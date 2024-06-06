from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import json
import os
from flask import Flask, send_file

app = Flask(__name__)

class PPTGenerator:
    def __init__(self, contents) -> None:
        self.contents = contents

    def addTitleSlide(self, prs: Presentation) -> None:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = self.contents['0']['mainTitle']
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(88, 24, 69)  # Red color

    def addTableOfContents(self, prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Contents"
        title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(88, 24, 69)  # Red color
        bullet_points = slide.shapes.placeholders[1]

        for key, value in self.contents.items():
            if 'title' in value:
                subheading = value["title"]
                text_frame = bullet_points.text_frame
                p = text_frame.add_paragraph()
                p.text = subheading
                p.level = 0
                p.runs[0].font.color.rgb = RGBColor(88, 24, 69)  # Blue color

    def addSlides(self, prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]
        items = list(self.contents.items())
        index = 0
        while index < len(items):
            current_item = items[index][1]
            if 'title' in current_item:
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = current_item['title']
                title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(88, 24, 69)  # Red color
                index += 1
                if index < len(items):
                    next_item = items[index][1]
                    if 'text' in next_item:
                        bullet_points = slide.shapes.placeholders[1]
                        text_frame = bullet_points.text_frame
                        p = text_frame.add_paragraph()
                        p.text = next_item['text']
                        p.level = 0
                        p.runs[0].font.color.rgb = RGBColor(0, 0, 255)  # Blue color
                        index += 1
            else:
                index += 1
                
    def removeFirstSlide(self, prs: Presentation) -> None:
        del prs.slides._sldIdLst[0]
        
        
@app.route("/generate")
def generate_ppt():
    # Load contents from the JSON file
    print("running")
    with open("./1.json", "r") as file:
        content = json.load(file)
    cwd = os.getcwd()
    newContent = {
        "Presentation Name": os.path.join(cwd, "Sample.pptx"),
    }
    newContent.update(content)

    MyPptGenerator = PPTGenerator(contents=newContent)
    prs = Presentation("./input/title.pptx")

    # Clear existing slides except the first one
    for i in range(len(prs.slides)-1, 0, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    MyPptGenerator.addTitleSlide(prs)
    MyPptGenerator.addTableOfContents(prs)
    MyPptGenerator.addSlides(prs)
    MyPptGenerator.removeFirstSlide(prs)

    return send_file(newContent["Presentation Name"], as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True, port=7000, host="0.0.0.0")